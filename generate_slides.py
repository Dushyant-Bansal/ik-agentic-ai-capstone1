"""Generate the capstone presentation slide deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x00, 0x7A, 0xCC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
MUTED = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x00, 0x96, 0x6E)
ORANGE = RGBColor(0xE8, 0x6C, 0x00)


def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _tf(shape, text, size=18, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf


def _add_bullet(tf, text, size=16, color=DARK_TEXT, level=0, bold=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.level = level
    p.space_before = Pt(4)


def _title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, DARK)
    _add_shape_bg(slide, Inches(0), Inches(2.5), Inches(10), Inches(0.06), ACCENT)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(8.4), Inches(1.5))
    _tf(tb, "AI-Powered Email Assistant", size=36, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)
    _add_bullet(tb.text_frame, "Multi-Agent Pipeline with LangGraph, Streamlit & Pydantic", size=20, color=RGBColor(0xBB, 0xBB, 0xBB))

    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(8.4), Inches(1.2))
    tf2 = _tf(tb2, "IK Agentic AI for SWEs - Capstone Project", size=18, color=RGBColor(0x99, 0x99, 0x99))
    _add_bullet(tf2, "Applied Agentic AI for Software Engineers", size=14, color=MUTED)


def _problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.7))
    _tf(tb, "Problem Statement & Business Opportunity", size=28, bold=True, color=DARK)

    _add_shape_bg(slide, Inches(0.6), Inches(1.2), Inches(8.8), Inches(0.04), ACCENT)

    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(4.0))
    tf = _tf(tb2, "The Problem", size=20, bold=True, color=ACCENT)
    _add_bullet(tf, "Professionals spend 15-20 min per email drafting", size=16, color=DARK_TEXT)
    _add_bullet(tf, "Inconsistent tone, formatting, and quality across teams", size=16, color=DARK_TEXT)
    _add_bullet(tf, "No context awareness between related emails", size=16, color=DARK_TEXT)
    _add_bullet(tf, "", size=10)
    _add_bullet(tf, "The Opportunity", size=20, bold=True, color=GREEN)
    _add_bullet(tf, "Reduce drafting time to under 2 minutes", size=16, color=DARK_TEXT)
    _add_bullet(tf, "Context-aware writing tailored to recipient and situation", size=16, color=DARK_TEXT)
    _add_bullet(tf, "Multi-purpose: outreach, follow-ups, apologies, internal updates", size=16, color=DARK_TEXT)
    _add_bullet(tf, "Consistent, on-brand communication at team/org level", size=16, color=DARK_TEXT)


def _agentic_principles_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.7))
    _tf(tb, "Agentic AI Principles Applied", size=28, bold=True, color=DARK)
    _add_shape_bg(slide, Inches(0.6), Inches(1.2), Inches(8.8), Inches(0.04), ACCENT)

    principles = [
        ("Modular Agent Decomposition", "Each task (parse, classify, tone, draft, personalize, review, route) is a dedicated agent with a single responsibility."),
        ("Structured State Management", "Pydantic-validated TypedDict flows through the graph; agents read and return typed partial updates."),
        ("Orchestrated Collaboration", "LangGraph wires agents into a directed graph with sequential edges and conditional retry loops."),
        ("Tool Use & Function Calling", "Agents use LLM structured output (with_structured_output) for reliable Pydantic-based parsing."),
        ("Memory & Personalization", "Conversation history persists across sessions; the Draft Writer uses recent turns as context."),
        ("Human-in-the-Loop", "Streamlit UI lets users edit drafts, override tone/intent, and manage their profile before export."),
    ]

    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.0))
    tf = tb2.text_frame
    tf.word_wrap = True
    first = True
    for title, desc in principles:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ACCENT
        p.space_before = Pt(8)
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = DARK_TEXT
        p2.space_before = Pt(2)


def _architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.7))
    _tf(tb, "System Architecture", size=28, bold=True, color=DARK)
    _add_shape_bg(slide, Inches(0.6), Inches(1.2), Inches(8.8), Inches(0.04), ACCENT)

    components = [
        ("Multi-Agent System", "LangGraph", "Modular task-specific agent orchestration"),
        ("Language Models", "GPT-4o-mini / Claude", "Language understanding and generation"),
        ("Email Template Engine", "LangChain + Function Calling", "Tone, length, and formatting control"),
        ("User Profile Store", "JSON (swappable to DB)", "Preferences and prior context"),
        ("Web Interface", "Streamlit", "Compose, edit, and export emails"),
        ("Memory Layer", "LangGraph + JSON", "Remembers user tone and previous drafts"),
        ("Control Plane", "config/mcp.yaml", "Model selection and routing"),
    ]

    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.0))
    tf = tb2.text_frame
    tf.word_wrap = True

    header = tf.paragraphs[0]
    header.text = f"{'Component':<28}{'Tech':<28}{'Purpose'}"
    header.font.size = Pt(12)
    header.font.bold = True
    header.font.color.rgb = ACCENT

    for comp, tech, purpose in components:
        p = tf.add_paragraph()
        p.text = f"{comp:<28}{tech:<28}{purpose}"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(4)


def _pipeline_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.7))
    _tf(tb, "Agent Pipeline", size=28, bold=True, color=DARK)
    _add_shape_bg(slide, Inches(0.6), Inches(1.2), Inches(8.8), Inches(0.04), ACCENT)

    agents = [
        ("1. Input Parser", "Validates prompt, extracts recipient, tone, constraints via LLM structured output"),
        ("2. Intent Detection", "Classifies intent: outreach, follow-up, apology, info request, internal update"),
        ("3. Tone Stylist", "Maps tone to prompt snippets; loads examples from tone_samples/"),
        ("4. Draft Writer", "Generates subject + body using tone context and conversation history"),
        ("5. Personalization", "Injects user name, company, signature; strips LLM placeholders"),
        ("6. Review & Validator", "Checks grammar, tone alignment, coherence; flags issues"),
        ("7. Router & Memory", "Logs to memory, decides retry (up to 2x) or finalizes draft"),
    ]

    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.0))
    tf = tb2.text_frame
    tf.word_wrap = True
    first = True
    for name, desc in agents:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = name
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ACCENT
        p.space_before = Pt(8)
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = DARK_TEXT
        p2.space_before = Pt(2)


def _langgraph_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.7))
    _tf(tb, "LangGraph Orchestration", size=28, bold=True, color=DARK)
    _add_shape_bg(slide, Inches(0.6), Inches(1.2), Inches(8.8), Inches(0.04), ACCENT)

    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.0))
    tf = _tf(tb2, "Why LangGraph?", size=20, bold=True, color=ACCENT)
    _add_bullet(tf, "Directed graph with typed state (TypedDict + Pydantic)", size=15, color=DARK_TEXT)
    _add_bullet(tf, "Conditional edges enable retry loops (Router -> Draft Writer)", size=15, color=DARK_TEXT)
    _add_bullet(tf, "In-memory checkpointer for debugging and state inspection", size=15, color=DARK_TEXT)
    _add_bullet(tf, "", size=10)
    _add_bullet(tf, "Graph Structure", size=20, bold=True, color=ACCENT)
    _add_bullet(tf, "7 nodes (one per agent class), 7 sequential edges", size=15, color=DARK_TEXT)
    _add_bullet(tf, "1 conditional edge: Router -> Draft Writer (if review fails, retry_count < 2)", size=15, color=DARK_TEXT)
    _add_bullet(tf, "1 conditional edge: Router -> END (if review passes or retries exhausted)", size=15, color=DARK_TEXT)
    _add_bullet(tf, "", size=10)
    _add_bullet(tf, "State Flow", size=20, bold=True, color=ACCENT)
    _add_bullet(tf, "Each node reads shared state, returns partial dict updates", size=15, color=DARK_TEXT)
    _add_bullet(tf, "State keys: parsed_input, intent, tone_context, draft, personalized_draft, review_result, errors, retry_count", size=14, color=MUTED)


def _memory_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.7))
    _tf(tb, "Memory & Personalization", size=28, bold=True, color=DARK)
    _add_shape_bg(slide, Inches(0.6), Inches(1.2), Inches(8.8), Inches(0.04), ACCENT)

    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.0))
    tf = _tf(tb2, "Conversation Memory", size=20, bold=True, color=ACCENT)
    _add_bullet(tf, "Each email generation is logged as a ConversationTurn (prompt, subject, body, intent, tone)", size=15, color=DARK_TEXT)
    _add_bullet(tf, "Last 10 turns stored per user in user_profiles.json", size=15, color=DARK_TEXT)
    _add_bullet(tf, "Draft Writer injects last 3 turns as context for style consistency", size=15, color=DARK_TEXT)
    _add_bullet(tf, "Users can clear history to start fresh via UI", size=15, color=DARK_TEXT)
    _add_bullet(tf, "", size=10)
    _add_bullet(tf, "User Profiles (Pydantic-backed)", size=20, bold=True, color=ACCENT)
    _add_bullet(tf, "Name, company, style preferences (signature, preferred/avoided phrases)", size=15, color=DARK_TEXT)
    _add_bullet(tf, "Prior draft summaries (last 20)", size=15, color=DARK_TEXT)
    _add_bullet(tf, "Personalization Agent uses profile to inject name/signature and replace placeholders", size=15, color=DARK_TEXT)
    _add_bullet(tf, "", size=10)
    _add_bullet(tf, "Placeholder Safety Net", size=20, bold=True, color=ACCENT)
    _add_bullet(tf, "Draft Writer prompt includes sender name + explicit no-placeholder instruction", size=15, color=DARK_TEXT)
    _add_bullet(tf, "Personalization Agent strips [Your Name], [Sender Name], etc. as fallback", size=15, color=DARK_TEXT)


def _eval_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.7))
    _tf(tb, "Evaluation & Testing", size=28, bold=True, color=DARK)
    _add_shape_bg(slide, Inches(0.6), Inches(1.2), Inches(8.8), Inches(0.04), ACCENT)

    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.0))
    tf = _tf(tb2, "Unit Tests (37 tests, no API calls)", size=20, bold=True, color=GREEN)
    _add_bullet(tf, "test_schemas.py: All Pydantic models, enums, validation, JSON round-trips", size=15, color=DARK_TEXT)
    _add_bullet(tf, "test_profile_store.py: Save/load/append/clear with temp JSON fixture", size=15, color=DARK_TEXT)
    _add_bullet(tf, "test_tone_stylist.py: Tone context generation for all 5 tones", size=15, color=DARK_TEXT)
    _add_bullet(tf, "test_personalization.py: Name/signature, placeholder stripping, company injection", size=15, color=DARK_TEXT)
    _add_bullet(tf, "", size=10)
    _add_bullet(tf, "LLM-as-a-Judge Evaluation (6 tests, uses API)", size=20, bold=True, color=ORANGE)
    _add_bullet(tf, "5 tone tests: generate email per tone, judge verifies match via structured output", size=15, color=DARK_TEXT)
    _add_bullet(tf, "1 consistency test: same prompt with formal vs casual, judge confirms tones differ", size=15, color=DARK_TEXT)
    _add_bullet(tf, "ToneJudgment schema: matches_requested_tone, detected_tone, confidence, reasoning", size=15, color=DARK_TEXT)
    _add_bullet(tf, "Auto-skips when OPENAI_API_KEY not set (CI-safe)", size=15, color=DARK_TEXT)


def _ui_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.7))
    _tf(tb, "Streamlit UI", size=28, bold=True, color=DARK)
    _add_shape_bg(slide, Inches(0.6), Inches(1.2), Inches(8.8), Inches(0.04), ACCENT)

    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(4.0), Inches(5.0))
    tf = _tf(tb2, "Sidebar Controls", size=20, bold=True, color=ACCENT)
    _add_bullet(tf, "Tone selector (5 options)", size=15, color=DARK_TEXT)
    _add_bullet(tf, "Intent override dropdown", size=15, color=DARK_TEXT)
    _add_bullet(tf, "Recipient field", size=15, color=DARK_TEXT)
    _add_bullet(tf, "User ID for personalization", size=15, color=DARK_TEXT)
    _add_bullet(tf, "Profile management (name, company)", size=15, color=DARK_TEXT)
    _add_bullet(tf, "Clear conversation history button", size=15, color=DARK_TEXT)

    tb3 = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.0), Inches(5.0))
    tf2 = _tf(tb3, "Main Area", size=20, bold=True, color=ACCENT)
    _add_bullet(tf2, "Prompt text area", size=15, color=DARK_TEXT)
    _add_bullet(tf2, "Generate Email button", size=15, color=DARK_TEXT)
    _add_bullet(tf2, "Editable subject + body preview", size=15, color=DARK_TEXT)
    _add_bullet(tf2, "Export as TXT download", size=15, color=DARK_TEXT)
    _add_bullet(tf2, "", size=10)
    _add_bullet(tf2, "Design Principles", size=20, bold=True, color=ACCENT)
    _add_bullet(tf2, "Human-in-the-loop editing", size=15, color=DARK_TEXT)
    _add_bullet(tf2, "Session state persistence", size=15, color=DARK_TEXT)
    _add_bullet(tf2, "Real-time feedback (errors, warnings)", size=15, color=DARK_TEXT)


def _design_decisions_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.7))
    _tf(tb, "Key Design Decisions", size=28, bold=True, color=DARK)
    _add_shape_bg(slide, Inches(0.6), Inches(1.2), Inches(8.8), Inches(0.04), ACCENT)

    decisions = [
        ("LangGraph over pure LangChain", "Built-in state management, conditional edges, retry loops, and checkpointing"),
        ("Pydantic for all agent I/O", "Structured validation, schema-driven prompts, type safety across the pipeline"),
        ("Class-based agents", "Consistent pattern across all 7 agents; each exposes run(state) -> dict"),
        ("JSON user profiles", "Simple, spec-compliant; easily swappable to MongoDB or Postgres"),
        ("LLM structured output", "with_structured_output() for reliable parsing and classification"),
        ("Conversation memory", "Last 10 turns persisted per user; last 3 injected as draft context"),
        ("LLM-as-a-Judge evaluation", "Automated tone verification using a second LLM call with Pydantic schema"),
        ("MCP config (mcp.yaml)", "Primary/fallback model routing; max_retries for review loop"),
    ]

    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.0))
    tf = tb2.text_frame
    tf.word_wrap = True
    first = True
    for title, desc in decisions:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = ACCENT
        p.space_before = Pt(6)
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = DARK_TEXT
        p2.space_before = Pt(1)


def _closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK)
    _add_shape_bg(slide, Inches(0), Inches(2.5), Inches(10), Inches(0.06), ACCENT)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(1.0))
    _tf(tb, "Thank You", size=36, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(8.4), Inches(3.0))
    tf = _tf(tb2, "Run it yourself:", size=18, color=RGBColor(0x99, 0x99, 0x99))
    _add_bullet(tf, "pip install -r requirements.txt", size=16, color=RGBColor(0xBB, 0xBB, 0xBB))
    _add_bullet(tf, "streamlit run email_assistant/src/ui/streamlit_app.py", size=16, color=RGBColor(0xBB, 0xBB, 0xBB))
    _add_bullet(tf, "", size=10)
    _add_bullet(tf, "Questions?", size=22, bold=True, color=WHITE)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _title_slide(prs)
    _problem_slide(prs)
    _agentic_principles_slide(prs)
    _architecture_slide(prs)
    _pipeline_slide(prs)
    _langgraph_slide(prs)
    _memory_slide(prs)
    _eval_slide(prs)
    _ui_slide(prs)
    _design_decisions_slide(prs)
    _closing_slide(prs)

    out = "AI_Email_Assistant_Presentation.pptx"
    prs.save(out)
    print(f"Saved {out} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
